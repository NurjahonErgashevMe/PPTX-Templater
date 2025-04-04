from pptx import Presentation
import re
from typing import Dict, List
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.util import Pt


class PowerPointTemplateUpdater:
    PLACEHOLDER_PATTERN = r"\{\{([^}]+)\}\}"

    def __init__(self, template_path: str, output_path: str):
        self.template_path = template_path
        self.output_path = output_path
        self.presentation = Presentation(template_path)

    @staticmethod
    def _find_placeholders(text: str) -> List[str]:
        return re.findall(PowerPointTemplateUpdater.PLACEHOLDER_PATTERN, text)

    def _process_shape(self, shape, replacements: Dict[str, str]):
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return

        original_text = shape.text_frame.text
        placeholders = self._find_placeholders(original_text)

        if not placeholders:
            return

        # Сохраняем ВСЕ параметры форматирования из первого run
        original_font = None
        if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
            original_run = shape.text_frame.paragraphs[0].runs[0]
            original_font = {
                'name': original_run.font.name,
                'size': original_run.font.size,
                'bold': original_run.font.bold,
                'italic': original_run.font.italic,
                'underline': original_run.font.underline,
                'color': self._get_actual_color(original_run.font)
            }

        # Замена плейсхолдеров
        new_text = original_text
        for placeholder in placeholders:
            key = placeholder.strip()
            new_text = new_text.replace(f"{{{{{key}}}}}", str(replacements.get(key, "")))

        # Полная очистка и восстановление
        shape.text_frame.clear()
        paragraph = shape.text_frame.add_paragraph()
        run = paragraph.add_run()
        run.text = new_text

        # Точное восстановление форматирования
        if original_font:
            self._apply_exact_font_style(run.font, original_font)

    @staticmethod
    def _get_actual_color(font):
        """Получает точное значение цвета, включая темы и RGB"""
        if not font.color:
            return None
            
        if font.color.type == MSO_COLOR_TYPE.RGB:
            return font.color.rgb
        elif font.color.type == MSO_COLOR_TYPE.SCHEME:
            return font.color.theme_color
        elif font.color.type == MSO_COLOR_TYPE.SYSTEM:
            return None  # Системные цвета не поддерживаем
        else:
            try:
                # Пытаемся получить RGB даже для нестандартных цветов
                return RGBColor(
                    font.color.rgb.r,
                    font.color.rgb.g,
                    font.color.rgb.b
                )
            except:
                return None

    @staticmethod
    def _apply_exact_font_style(target_font, source_style):
        """Точное применение всех стилей"""
        if source_style['name']:
            target_font.name = source_style['name']
        if source_style['size']:
            target_font.size = source_style['size']
        if source_style['bold'] is not None:
            target_font.bold = source_style['bold']
        if source_style['italic'] is not None:
            target_font.italic = source_style['italic']
        if source_style['underline'] is not None:
            target_font.underline = source_style['underline']
        
        # Особое внимание для цвета
        if source_style['color']:
            try:
                if isinstance(source_style['color'], RGBColor):
                    target_font.color.rgb = source_style['color']
                else:
                    # Для цветов темы
                    target_font.color.theme_color = source_style['color']
            except Exception as e:
                print(f"Не удалось применить цвет: {str(e)}")

    def update_template(self, data: Dict[str, str]):
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                self._process_shape(shape, data)
        
        self.presentation.save(self.output_path)


if __name__ == "__main__":
    try:
        data = {
            "theme": "Школа",
            "description": "описание школы* lorem ipsum dolor..."
        }
        
        updater = PowerPointTemplateUpdater(
            template_path="template.pptx",
            output_path="output.pptx"
        )
        updater.update_template(data)
        print("Презентация обновлена с полным сохранением форматирования!")
    except Exception as e:
        print(f"Ошибка: {str(e)}")
        raise